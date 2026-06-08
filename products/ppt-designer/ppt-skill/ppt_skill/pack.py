"""Local packaging: assemble a deck's images into .pptx and into an HTML viewer."""

import json
import os
import re
import sys
from pathlib import Path

TEMPLATES_DIR = Path(__file__).resolve().parent.parent / "templates"
VIEWER_TEMPLATE = TEMPLATES_DIR / "viewer.html"

_PNG_MAGIC = b"\x89PNG\r\n\x1a\n"
_SLIDE_PATTERN = re.compile(r"^slide-(\d+)\.png$")
_MIN_PNG_BYTES = 10_000  # anything smaller is almost certainly corrupt/placeholder
_MAX_RATIO_DEVIATION = 0.08  # warn if aspect ratio differs from 16:9 by more than this


def _validate_png(path: Path) -> None:
    """Check magic bytes, minimum file size, and aspect ratio. Raises on hard failures."""
    try:
        data = path.read_bytes()
    except OSError as e:
        raise ValueError(f"Cannot read {path.name}: {e}") from e
    if len(data) < 24 or data[:8] != _PNG_MAGIC:
        raise ValueError(f"Not a valid PNG (bad magic bytes): {path.name}")
    if len(data) < _MIN_PNG_BYTES:
        raise ValueError(
            f"PNG too small ({len(data)} bytes) — likely corrupt or placeholder: {path.name}"
        )
    w = int.from_bytes(data[16:20], "big")
    h = int.from_bytes(data[20:24], "big")
    if h > 0 and abs(w / h - 16 / 9) > _MAX_RATIO_DEVIATION:
        print(
            f"  [warn] {path.name}: aspect ratio {w}x{h} ({w/h:.3f}) deviates from 16:9 — slide may have letterboxing",
            file=sys.stderr,
        )


def to_pptx(images: list[Path], out_path: Path, notes: dict[int, str] | None = None) -> None:
    """Build a 16:9 .pptx — each image becomes one full-bleed slide.

    notes: optional mapping of slot number → speaker-note text.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError:
        print("Error: python-pptx not installed. Run: pip install -r requirements.txt", file=sys.stderr)
        sys.exit(1)

    if not images:
        raise ValueError("to_pptx: no images provided")

    prs = Presentation()
    # 16:9 canvas: 13.333 x 7.5 inches
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    blank_layout = prs.slide_layouts[6]

    for img in images:
        _validate_png(img)
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
        if notes:
            m = _SLIDE_PATTERN.match(img.name)
            slot = int(m.group(1)) if m else None
            if slot and (note := notes.get(slot)):
                slide.notes_slide.notes_text_frame.text = note

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"[pptx] {len(images)} slides -> {out_path}")


def to_html(images: list[Path], out_path: Path, title: str) -> None:
    """Build a self-contained HTML viewer for a list of images."""
    if not images:
        raise ValueError("to_html: no images provided")
    if not VIEWER_TEMPLATE.exists():
        raise FileNotFoundError(f"viewer template missing: {VIEWER_TEMPLATE}")

    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_dir = out_path.parent.resolve()

    rel_paths = []
    for img in images:
        try:
            rel = os.path.relpath(img.resolve(), out_dir)
        except ValueError:
            rel = str(img.resolve())
        rel_paths.append(rel)

    html = (VIEWER_TEMPLATE.read_text(encoding="utf-8")
            .replace("__TITLE__", title)
            .replace("__COUNT__", str(len(images)))
            .replace("__SLIDES__", json.dumps(rel_paths)))

    out_path.write_text(html, encoding="utf-8")
    print(f"[html] {len(images)} slides -> {out_path}")
    print(f"[html] open: file://{out_path.resolve()}")
